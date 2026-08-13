"""Generador PPTX fiel a Presentacion_Mantenimiento_GOS / Abastecimiento."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from gos.services.presentacion_catalog import (
    COLOR_DARK,
    COLOR_DARK_ALT,
    COLOR_FOOTER,
    COLOR_GOLD,
    COLOR_GRAY,
    COLOR_GREEN,
    COLOR_LIGHT,
    COLOR_MUTED,
)

SLIDE_W = 12191695
SLIDE_H = 6858000
M = 502920
TOP_EYEBROW = 384048
TOP_TITLE = 658368
FOOTER_Y = 6455664

RGB = {
    "gold": RGBColor.from_string(COLOR_GOLD),
    "green": RGBColor.from_string(COLOR_GREEN),
    "dark": RGBColor.from_string(COLOR_DARK),
    "dark2": RGBColor.from_string("2F3438"),
    "dark_alt": RGBColor.from_string(COLOR_DARK_ALT),
    "gray": RGBColor.from_string(COLOR_GRAY),
    "muted": RGBColor.from_string(COLOR_MUTED),
    "footer": RGBColor.from_string(COLOR_FOOTER),
    "light": RGBColor.from_string(COLOR_LIGHT),
    "chip": RGBColor.from_string("FFF6DA"),
    "white": RGBColor(255, 255, 255),
    "bar_track": RGBColor.from_string("DCDDDF"),
}

LOGO_PATH = Path(__file__).resolve().parent.parent / "static" / "img" / "gos-logo.png"


def build_presentation(
    module: dict,
    slides: list[dict],
    *,
    overview_kpis: list[dict] | None = None,
    chips: list[str] | None = None,
    periodo: str = "",
    empresa: str = "GREEN OIL SERVICES",
    dark_bar=None,
    circuit=None,
) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    n = 1
    _slide_cover(prs, blank, module, empresa, n)
    n += 1
    _slide_statement(prs, blank, module, n)
    n += 1
    _slide_flujo(prs, blank, module, chips or [], n)
    n += 1
    _slide_overview(prs, blank, module, overview_kpis or [], periodo, n)
    n += 1
    for spec in slides:
        n = _render_content_slide(prs, blank, spec, n)
    _slide_cierre(prs, blank, module, slides, empresa, n)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── primitives ──────────────────────────────────────────────────────────────


def _run(p, text, *, size, bold, color):
    r = p.add_run()
    r.text = "" if text is None else str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return r


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _lines(slide, left, top, width, height, lines, *, align=PP_ALIGN.LEFT):
    tf = _textbox(slide, left, top, width, height)
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        _run(p, text, size=size, bold=bold, color=color)


def _mixed(slide, left, top, width, height, parts, *, size=26, align=PP_ALIGN.LEFT):
    tf = _textbox(slide, left, top, width, height)
    p = tf.paragraphs[0]
    p.alignment = align
    for text, color, bold in parts:
        _run(p, text, size=size, bold=bold, color=color)


def _rect(slide, left, top, width, height, color, *, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1.25)
    return sh


def _oval(slide, left, top, width, height, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _top_bar(slide):
    _rect(slide, 0, 0, 8778021, 64008, RGB["gold"])
    _rect(slide, 8778021, 0, 1584920, 64008, RGB["gray"])
    _rect(slide, 10362941, 0, 1828754, 64008, RGB["dark"])
    _rect(slide, 9814255, -1737360, 2377440, 2377440, RGB["gray"])


def _footer(slide, text, number, *, dark=False):
    color = RGB["footer"] if dark else RGB["muted"]
    _lines(slide, M, FOOTER_Y, 7315200, 274320, [(text or "", 9, True, color)])
    _lines(slide, 10774375, FOOTER_Y, 914400, 274320, [(f"{number:02d}", 9, False, color)], align=PP_ALIGN.RIGHT)


def _title_block(slide, eyebrow, title, highlight=None, *, size=24):
    _lines(slide, M, TOP_EYEBROW, 10515600, 292608, [((eyebrow or "").upper(), 12, True, RGB["green"])])
    title = title or ""
    if highlight and highlight in title:
        before, after = title.split(highlight, 1)
        parts = [(before, RGB["dark"], True), (highlight, RGB["gold"], True), (after, RGB["dark"], True)]
    else:
        parts = [(title, RGB["dark"], True)]
    _mixed(slide, M, TOP_TITLE, 11155680, 800000, parts, size=size)


def _kpi_row(slide, kpis, *, top=1783080):
    tones = [
        (RGB["dark2"], RGB["white"]),
        (RGB["gold"], RGB["dark"]),
        (RGB["green"], RGB["white"]),
        (RGB["white"], RGB["dark"]),
    ]
    w, h = 2604440, 1234440
    gap = (11185855 - 4 * w) // 3
    for i, kpi in enumerate((kpis or [])[:4]):
        left = M + i * (w + gap)
        fill, tc = tones[i]
        _rect(slide, left, top, w, h, fill, line=RGB["dark"] if i == 3 else None)
        _lines(slide, left + 180000, top + 120000, w - 360000, 600000, [(str(kpi.get("value") or "—"), 26, True, tc)])
        _lines(slide, left + 180000, top + 780000, w - 360000, 400000, [(str(kpi.get("label") or ""), 11, None, tc)])


def _white_slide(prs, blank):
    slide = prs.slides.add_slide(blank)
    _bg(slide, RGB["white"])
    _top_bar(slide)
    return slide


# ── fixed intro/outro ───────────────────────────────────────────────────────


def _slide_cover(prs, blank, module, empresa, number):
    slide = prs.slides.add_slide(blank)
    _bg(slide, RGB["dark"])
    _rect(slide, 9814255, -1737360, 2377440, 2377440, RGB["dark_alt"])
    _oval(slide, 8869680, -1463040, 5029200, 5029200, RGB["gold"])
    _oval(slide, 10332720, 2926080, 3840480, 3840480, RGB["green"])
    _lines(slide, M, 566928, 9144000, 320040, [(module["sistema"], 13, True, RGB["gold"])])
    _mixed(
        slide,
        M,
        2331720,
        10607040,
        1920240,
        [(module["tagline"] + "\n", RGB["white"], True), (module["tagline_accent"], RGB["gold"], True)],
        size=46,
    )
    pilares = module["pilares"][0] if module.get("pilares") else ""
    if isinstance(module.get("pilares"), list) and len(module["pilares"]) > 1:
        pilares = "  ·  ".join(module["pilares"])
    _lines(slide, M, 4224528, 10058400, 365760, [(pilares, 16, None, RGB["footer"])])
    _rect(slide, M, 4828032, 5486400, 502920, RGB["gold"])
    _lines(slide, 685800, 4828032, 5212080, 502920, [(module["callout"], 13.5, True, RGB["dark"])])
    _footer(slide, empresa, number, dark=True)


def _slide_statement(prs, blank, module, number):
    slide = _white_slide(prs, blank)
    _lines(slide, M, 384048, 5486400, 274320, [("GREEN OIL SERVICES", 12, True, RGB["green"])])
    if LOGO_PATH.exists():
        try:
            slide.shapes.add_picture(str(LOGO_PATH), 5181448, 1417320, height=1527048)
        except Exception:
            pass
    _lines(slide, 1097280, 3703320, 9997135, 914400, [(module["statement"], 26, True, RGB["dark"])], align=PP_ALIGN.CENTER)
    _lines(
        slide,
        1097280,
        4617720,
        9997135,
        365760,
        [("  ·  ".join(module.get("statement_pilares") or []), 13.5, True, RGB["green"])],
        align=PP_ALIGN.CENTER,
    )
    _footer(slide, module["sistema"], number)


def _slide_flujo(prs, blank, module, chips, number):
    slide = _white_slide(prs, blank)
    _title_block(
        slide,
        "BASE DE DATOS Y HERRAMIENTAS DE GESTIÓN",
        "Una única base transforma reportes en decisiones",
        "decisiones",
        size=27,
    )
    _rect(slide, M, 1783080, 3017520, 2240280, RGB["light"])
    _lines(slide, M, 1920240, 3017520, 900000, [("⚙", 36, True, RGB["dark"])], align=PP_ALIGN.CENTER)
    _lines(
        slide,
        M,
        2900000,
        3017520,
        700000,
        [("MÓDULO", 12, True, RGB["muted"]), (module["label"].upper(), 16, True, RGB["dark"])],
        align=PP_ALIGN.CENTER,
    )
    _lines(slide, M, 3566160, 3017520, 320040, [("GOS Objetivos", 12, None, RGB["muted"])], align=PP_ALIGN.CENTER)

    flujo = (module.get("flujo") or [])[:4]
    box_w, box_h, arrow_w = 1390574, 1051560, 347472
    x, y = 3931920, 2011680
    for i, step in enumerate(flujo):
        last = i == len(flujo) - 1
        fill = RGB["green"] if last else RGB["white"]
        tc = RGB["white"] if last else RGB["dark"]
        _rect(slide, x, y, box_w, box_h, fill, line=None if last else RGB["gray"])
        _lines(slide, x + 100000, y + 180000, box_w - 200000, 400000, [(step.get("title", ""), 11, True, tc)], align=PP_ALIGN.CENTER)
        _lines(
            slide,
            x + 100000,
            y + 600000,
            box_w - 200000,
            350000,
            [(step.get("subtitle", ""), 10, None, tc if last else RGB["muted"])],
            align=PP_ALIGN.CENTER,
        )
        x += box_w
        if not last:
            _lines(slide, x, y, arrow_w, box_h, [("→", 18, True, RGB["gold"])], align=PP_ALIGN.CENTER)
            x += arrow_w

    chip_x = M
    for label in (chips or [])[:6]:
        w = max(827532, min(1444752, int(120000 + len(label) * 90000)))
        _rect(slide, chip_x, 4526280, w, 457200, RGB["chip"])
        _lines(slide, chip_x, 4526280, w, 457200, [(label, 12, True, RGB["dark"])], align=PP_ALIGN.CENTER)
        chip_x += w + 201168
    _footer(slide, "TRAZABILIDAD DE PUNTA A PUNTA", number)


def _slide_overview(prs, blank, module, kpis, periodo, number):
    slide = _white_slide(prs, blank)
    eyebrow = f"ESTADO GENERAL DEL {module['label'].upper()}"
    if periodo:
        eyebrow = f"{eyebrow} · {periodo}"
    _title_block(slide, eyebrow, module["overview_title"], module.get("overview_highlight"), size=26)
    if not kpis:
        kpis = [
            {"value": str(len(module.get("submodulos") or [])), "label": "submódulos"},
            {"value": str(len(module.get("proceso") or [])), "label": "etapas"},
            {"value": "GOS", "label": "plataforma"},
            {"value": periodo or "—", "label": "período"},
        ]
    _kpi_row(slide, kpis)
    proceso = module.get("proceso") or []
    if proceso:
        bar_y, bar_h = 3429000, 566928
        _rect(slide, M, bar_y, 11185855, bar_h, RGB["light"])
        n = len(proceso)
        slot = 11185855 // max(n * 2 - 1, 1)
        x = M
        for i, step in enumerate(proceso):
            _lines(slide, x, bar_y, slot, bar_h, [(step, 11, True, RGB["dark"])], align=PP_ALIGN.CENTER)
            x += slot
            if i < n - 1:
                _lines(slide, x, bar_y, slot, bar_h, [("→", 14, True, RGB["gold"])], align=PP_ALIGN.CENTER)
                x += slot
    _lines(slide, M, 4224528, 11185855, 365760, [(module["statement"], 14, None, RGB["muted"])])
    _footer(slide, module.get("overview_footer") or "VISIÓN INTEGRAL", number)


def _slide_cierre(prs, blank, module, slides, empresa, number):
    slide = prs.slides.add_slide(blank)
    _bg(slide, RGB["dark"])
    _rect(slide, 9814255, -1737360, 2377440, 2377440, RGB["dark_alt"])
    _oval(slide, 8869680, -1463040, 5029200, 5029200, RGB["gold"])
    _oval(slide, 10332720, 2926080, 3840480, 3840480, RGB["green"])
    _lines(slide, M, 566928, 9144000, 320040, [(module["sistema"], 13, True, RGB["gold"])])
    _mixed(
        slide,
        M,
        2200000,
        10607040,
        1600000,
        [("Compromiso de gestión\n", RGB["white"], True), ("datos reales → trazabilidad → acción", RGB["gold"], True)],
        size=32,
    )
    labels = []
    for s in slides:
        lab = s.get("footer") or s.get("label") or s.get("eyebrow")
        if lab and lab not in labels:
            labels.append(lab)
    _lines(slide, M, 4200000, 10000000, 400000, [("  ·  ".join(labels[:6]), 13, None, RGB["footer"])])
    _rect(slide, M, 4828032, 5486400, 502920, RGB["gold"])
    _lines(slide, 685800, 4828032, 5212080, 502920, [("Generado desde GOS Objetivos", 13.5, True, RGB["dark"])])
    _footer(slide, empresa, number, dark=True)


# ── content layouts ─────────────────────────────────────────────────────────


def _render_content_slide(prs, blank, spec: dict, number: int) -> int:
    layout = (spec.get("layout") or "bullets").lower()
    handlers = {
        "circuit": _layout_circuit,
        "kpi_bars": _layout_kpi_bars,
        "kpi_rank": _layout_kpi_rank,
        "kpi_rows": _layout_kpi_rows,
        "ranking": _layout_ranking,
        "ranking_spotlight": _layout_ranking_spotlight,
        "team_vtv": _layout_team_vtv,
        "category_board": _layout_category_board,
        "matrix": _layout_matrix,
        "next_steps": _layout_next_steps,
        "team_roster": _layout_team_roster,
        "bullets": _layout_bullets,
    }
    fn = handlers.get(layout, _layout_bullets)
    slide = _white_slide(prs, blank)
    _title_block(slide, spec.get("eyebrow") or "", spec.get("title") or "", spec.get("highlight"), size=22)
    fn(slide, spec)
    if spec.get("note"):
        _lines(slide, M, 5989320, 11185855, 300000, [(spec["note"], 9, None, RGB["muted"])])
    _footer(slide, spec.get("footer") or "", number)
    return number + 1


def _draw_ranking(slide, title, items, *, left, top, width, limit=8):
    _lines(slide, left, top, width, 280000, [(title, 11, True, RGB["muted"])])
    y = top + 320000
    for i, it in enumerate((items or [])[:limit]):
        _lines(slide, left, y, 450000, 280000, [(str(it.get("rank") or f"{i+1:02d}"), 12, True, RGB["gold"])])
        _lines(slide, left + 480000, y, width - 2100000, 280000, [(str(it.get("label") or "—"), 12, True, RGB["dark"])])
        _lines(slide, left + width - 1500000, y, 1450000, 280000, [(str(it.get("value") or "—"), 12, True, RGB["dark"])], align=PP_ALIGN.RIGHT)
        track = width - 480000
        _rect(slide, left + 480000, y + 270000, track, 45000, RGB["bar_track"])
        fill_w = int(track * max(0.1, 1 - i * 0.07))
        fill = RGB["green"] if i == 0 else RGB["gold"] if i == 1 else RGB["dark"]
        _rect(slide, left + 480000, y + 270000, fill_w, 45000, fill)
        y += 390000


def _draw_bars(slide, title, items, *, left, top, width, limit=8):
    _lines(slide, left, top, width, 280000, [(title, 11, True, RGB["muted"])])
    y = top + 360000
    for i, it in enumerate((items or [])[:limit]):
        pct = max(0.0, min(1.0, float(it.get("pct") or 0)))
        _lines(slide, left, y, width, 220000, [(str(it.get("label") or "—"), 11, True, RGB["dark"])])
        _rect(slide, left, y + 230000, width, 130000, RGB["bar_track"])
        fill = RGB["green"] if i == 0 else RGB["gold"] if i == 1 else RGB["dark"]
        _rect(slide, left, y + 230000, max(60000, int(width * pct)), 130000, fill)
        _lines(slide, left, y + 230000, width - 80000, 130000, [(str(it.get("value") or ""), 10, True, RGB["white"] if pct > 0.2 else RGB["dark"])], align=PP_ALIGN.RIGHT)
        y += 450000


def _layout_circuit(slide, spec):
    cards = spec.get("cards") or []
    n = min(len(cards), 4) or 1
    w = 2604440
    gap = (11185855 - n * w) // max(n - 1, 1) if n > 1 else 0
    accents = {"dark": RGB["dark"], "green": RGB["green"], "gold": RGB["gold"]}
    top = 1700000
    for i, card in enumerate(cards[:n]):
        left = M + i * (w + gap)
        accent = accents.get(card.get("accent"), RGB["dark"])
        _rect(slide, left, top, w, 54864, accent)
        _rect(slide, left, top + 54864, w, 2200000, RGB["light"])
        _oval(slide, left + 201168, top + 274320, 502920, 502920, accent)
        _lines(slide, left + 201168, top + 340000, 502920, 400000, [(card.get("letter") or "·", 16, True, RGB["white"])], align=PP_ALIGN.CENTER)
        _lines(slide, left + 201168, top + 900000, w - 402336, 450000, [(card.get("title") or "", 13, True, RGB["dark"])])
        _lines(slide, left + 201168, top + 1400000, w - 402336, 700000, [(card.get("text") or "", 11, None, RGB["muted"])])
    bar = spec.get("dark_bar") or {}
    if bar:
        _rect(slide, M, 4434840, 11185855, 777240, RGB["dark"])
        _lines(slide, 731520, 4434840, 4572000, 777240, [(bar.get("left") or "", 12, True, RGB["gold"])])
        metrics = bar.get("metrics") or []
        x = 5349240
        for m in metrics[:3]:
            _lines(slide, x, 4507992, 2011680, 384048, [(str(m.get("value") or ""), 18, True, RGB["white"])])
            _lines(slide, x, 4892040, 2011680, 274320, [(str(m.get("label") or ""), 10, None, RGB["footer"])])
            x += 2011680


def _layout_kpi_bars(slide, spec):
    if spec.get("kpis"):
        _kpi_row(slide, spec["kpis"], top=1550000)
    _draw_bars(slide, spec.get("bars_title") or "DISTRIBUCIÓN", spec.get("bars") or [], left=M, top=3000000, width=11185855, limit=6)


def _layout_kpi_rank(slide, spec):
    if spec.get("kpis"):
        _kpi_row(slide, spec["kpis"], top=1500000)
    _draw_ranking(slide, spec.get("ranking_title") or "RANKING", spec.get("ranking") or [], left=M, top=2900000, width=5400000, limit=6)
    if spec.get("bars"):
        _draw_bars(slide, spec.get("bars_title") or "SERIE", spec["bars"], left=6200000, top=2900000, width=5400000, limit=6)


def _layout_kpi_rows(slide, spec):
    if spec.get("kpis"):
        _kpi_row(slide, spec["kpis"], top=1550000)
    rows = spec.get("rows") or []
    y = 3000000
    if spec.get("rows_title"):
        _lines(slide, M, y, 11185855, 280000, [(spec["rows_title"], 11, True, RGB["muted"])])
        y += 320000
    for row in rows[:7]:
        _rect(slide, M, y, 11185855, 480000, RGB["white"], line=RGB["gray"])
        _lines(slide, M + 150000, y + 80000, 3200000, 350000, [(str(row.get("a") or ""), 13, True, RGB["dark"])])
        _lines(slide, M + 4000000, y + 80000, 3200000, 350000, [(str(row.get("b") or ""), 12, None, RGB["muted"])])
        _lines(slide, M + 8500000, y + 80000, 2500000, 350000, [(str(row.get("c") or ""), 12, True, RGB["dark"])], align=PP_ALIGN.RIGHT)
        y += 520000


def _layout_ranking(slide, spec):
    _draw_ranking(slide, spec.get("ranking_title") or "RANKING", spec.get("ranking") or [], left=M, top=1600000, width=11185855, limit=10)


def _layout_ranking_spotlight(slide, spec):
    if spec.get("kpis"):
        _kpi_row(slide, spec["kpis"], top=1450000)
        top = 2850000
    else:
        top = 1600000
    _draw_ranking(slide, spec.get("ranking_title") or "TOP 10", spec.get("ranking") or [], left=M, top=top, width=7000000, limit=8)
    spot = spec.get("spotlight") or {}
    _rect(slide, 7800000, top, 3800000, 1600000, RGB["dark"])
    _lines(slide, 8000000, top + 250000, 3400000, 700000, [(str(spot.get("value") or "—"), 36, True, RGB["white"])])
    _lines(slide, 8000000, top + 1000000, 3400000, 500000, [(str(spot.get("text") or ""), 12, None, RGB["footer"])])


def _layout_team_vtv(slide, spec):
    # left oficios
    _rect(slide, M, 1550000, 3520440, 4200000, RGB["light"])
    _lines(slide, 685800, 1680000, 3154680, 300000, [("Distribución por oficio", 12, True, RGB["dark"])])
    y = 2100000
    oficios = spec.get("oficios") or []
    max_v = max((o.get("value") or 0) for o in oficios) or 1
    for o in oficios[:7]:
        _lines(slide, 685800, y, 2500000, 220000, [(str(o.get("label") or ""), 11, True, RGB["dark"])])
        bw = int(2500000 * ((o.get("value") or 0) / max_v))
        _rect(slide, 685800, y + 220000, 2697480, 120000, RGB["bar_track"])
        _rect(slide, 685800, y + 220000, max(80000, bw), 120000, RGB["green"])
        _lines(slide, 3429000, y, 400000, 250000, [(str(o.get("value") or 0), 12, True, RGB["dark"])])
        y += 480000

    # middle bases
    _rect(slide, 4315968, 1550000, 3520440, 4200000, RGB["light"])
    _lines(slide, 4498848, 1680000, 3154680, 300000, [("Distribución por base", 12, True, RGB["dark"])])
    y = 2200000
    for b in (spec.get("bases") or [])[:5]:
        _lines(slide, 4498848, y, 3154680, 250000, [(str(b.get("label") or ""), 12, True, RGB["dark"])])
        _rect(slide, 4498848, y + 280000, 3154680, 200000, RGB["bar_track"])
        _rect(slide, 4498848, y + 280000, max(100000, int(3154680 * float(b.get("pct") or 0))), 200000, RGB["dark"])
        _lines(slide, 4498848, y + 520000, 3154680, 220000, [(str(b.get("text") or ""), 11, None, RGB["muted"])])
        y += 900000
    _rect(slide, 4498848, 4800000, 3154680, 777240, RGB["dark"])
    _lines(slide, 4636008, 4880000, 1097280, 500000, [(str(spec.get("team_total") or "—"), 28, True, RGB["white"])])
    _lines(slide, 5550408, 4950000, 2000000, 500000, [("personas en el equipo\nde mantenimiento", 11, None, RGB["footer"])])

    # right VTV
    _rect(slide, 8129016, 1550000, 3520440, 4200000, RGB["light"])
    _lines(slide, 8311896, 1680000, 3154680, 300000, [("VTV — control legal", 12, True, RGB["dark"])])
    vk = (spec.get("vtv_kpis") or [{}])[0]
    _lines(slide, 8311896, 2100000, 3154680, 600000, [(str(vk.get("value") or "—"), 36, True, RGB["dark"])])
    _lines(slide, 8311896, 2750000, 3154680, 350000, [(str(vk.get("label") or ""), 11, None, RGB["muted"])])
    y = 3300000
    for row in (spec.get("vtv_rows") or [])[:4]:
        _rect(slide, 8311896, y, 3154680, 480000, RGB["white"], line=RGB["gray"])
        _lines(slide, 8449056, y + 70000, 1200000, 350000, [(str(row.get("a") or ""), 11, True, RGB["dark"])])
        _lines(slide, 9600000, y + 70000, 1000000, 350000, [(str(row.get("b") or ""), 10, None, RGB["muted"])])
        _lines(slide, 10500000, y + 70000, 900000, 350000, [(str(row.get("c") or ""), 10, True, RGB["dark"])], align=PP_ALIGN.RIGHT)
        y += 520000


def _layout_category_board(slide, spec):
    cats = spec.get("categories") or []
    n = min(len(cats), 8)
    cols = 4
    w, h = 2604440, 1600000
    gap_x = (11185855 - cols * w) // (cols - 1) if cols > 1 else 0
    for i, cat in enumerate(cats[:n]):
        row, col = divmod(i, cols)
        left = M + col * (w + gap_x)
        top = 1600000 + row * (h + 200000)
        accent = RGB["gold"] if i % 2 == 0 else RGB["green"]
        _rect(slide, left, top, w, 50000, accent)
        _rect(slide, left, top + 50000, w, h - 50000, RGB["light"])
        _lines(slide, left + 150000, top + 180000, w - 300000, 350000, [(str(cat.get("title") or ""), 13, True, RGB["dark"])])
        _lines(slide, left + 150000, top + 550000, w - 300000, 500000, [(str(cat.get("value") or "—"), 28, True, RGB["dark"])])
        _lines(slide, left + 150000, top + 1100000, w - 300000, 350000, [(str(cat.get("sub") or ""), 11, None, RGB["muted"])])


def _layout_matrix(slide, spec):
    headers = spec.get("headers") or []
    rows = spec.get("rows") or []
    if not headers:
        return
    cols = min(len(headers), 12)
    left0 = M + 1800000
    cell_w = (11185855 - 1800000) // max(cols, 1)
    y = 1600000
    _rect(slide, M, y, 11185855, 500000, RGB["light"])
    _lines(slide, M + 100000, y, 1600000, 500000, [("Categoría", 11, True, RGB["dark"])])
    for i, h in enumerate(headers[:cols]):
        _lines(slide, left0 + i * cell_w, y, cell_w, 500000, [(str(h), 10, True, RGB["dark"])], align=PP_ALIGN.CENTER)
    y += 520000
    for r, row in enumerate(rows[:6]):
        bg = RGB["white"] if r % 2 == 0 else RGB["light"]
        _rect(slide, M, y, 11185855, 450000, bg)
        _lines(slide, M + 100000, y, 1600000, 450000, [(str(row.get("label") or ""), 11, True, RGB["dark"])])
        cells = row.get("cells") or []
        for i, val in enumerate(cells[:cols]):
            _lines(slide, left0 + i * cell_w, y, cell_w, 450000, [(str(val), 11, None, RGB["dark"])], align=PP_ALIGN.CENTER)
        y += 460000
    # callouts
    x = M
    for c in (spec.get("callouts") or [])[:3]:
        _rect(slide, x, 5000000, 3600000, 900000, RGB["light"])
        _lines(slide, x + 150000, 5100000, 3300000, 300000, [(str(c.get("title") or ""), 14, True, RGB["gold"])])
        _lines(slide, x + 150000, 5450000, 3300000, 350000, [(str(c.get("text") or ""), 12, None, RGB["dark"])])
        x += 3800000


def _layout_next_steps(slide, spec):
    steps = spec.get("steps") or []
    for i, step in enumerate(steps[:6]):
        col = i % 3
        row = i // 3
        left = M + col * 3800000
        top = 1600000 + row * 2000000
        _rect(slide, left, top, 3600000, 1800000, RGB["light"])
        _lines(slide, left + 150000, top + 150000, 1000000, 400000, [(str(step.get("n") or f"{i+1:02d}"), 22, True, RGB["gold"])])
        _lines(slide, left + 150000, top + 600000, 3300000, 450000, [(str(step.get("title") or ""), 14, True, RGB["dark"])])
        _lines(slide, left + 150000, top + 1100000, 3300000, 550000, [(str(step.get("text") or ""), 12, None, RGB["muted"])])


def _layout_team_roster(slide, spec):
    people = spec.get("people") or []
    cols = 4
    w, h = 2604440, 900000
    gap = (11185855 - cols * w) // (cols - 1) if cols > 1 else 0
    for i, p in enumerate(people[:16]):
        row, col = divmod(i, cols)
        left = M + col * (w + gap)
        top = 1550000 + row * (h + 120000)
        _rect(slide, left, top, w, h, RGB["light"])
        _lines(slide, left + 120000, top + 150000, w - 240000, 350000, [(str(p.get("nombre") or "—"), 12, True, RGB["dark"])])
        _lines(slide, left + 120000, top + 480000, w - 240000, 300000, [(str(p.get("rol") or ""), 11, None, RGB["muted"])])
        _lines(slide, left + 120000, top + 700000, w - 240000, 150000, [(str(p.get("base") or ""), 10, True, RGB["green"])])


def _layout_bullets(slide, spec):
    if spec.get("kpis"):
        _kpi_row(slide, spec["kpis"], top=1500000)
        y = 3000000
    else:
        y = 1700000
    for b in (spec.get("bullets") or [])[:6]:
        _rect(slide, M, y + 80000, 110000, 110000, RGB["gold"])
        _lines(slide, M + 250000, y, 10500000, 500000, [(b, 15, None, RGB["dark"])])
        y += 580000
